"""
Generate Vocaband sales collateral: two A4 handouts + one full presentation,
in Hebrew and Arabic.

Outputs per language (lang = HE | AR):
- docs/sales/teacher-handout-{lang}.pptx + .pdf       (A4 portrait, single page)
- docs/sales/manager-handout-{lang}.pptx + .pdf       (A4 portrait, single page)
- docs/sales/vocaband-presentation-{lang}.pptx + .pdf (16:9 deck)

Re-run after editing this file to regenerate everything:
    python3 scripts/build-sales-collateral.py
"""

from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "docs" / "sales"
SHOTS = ROOT / "docs" / "screenshots"

OUT_DIR.mkdir(parents=True, exist_ok=True)


# --- Brand palette (Vocaband: indigo -> violet -> fuchsia + amber accent) ---
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
FUCHSIA = RGBColor(0xD9, 0x46, 0xEF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
INK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
BG_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- Language strings ----------

STRINGS = {
    "he": {
        "tag": "HE",
        "font": "Noto Sans Hebrew",

        # Page chrome / section labels
        "sec_problem": "הבעיה",
        "sec_basis": "על מה זה בנוי",
        "sec_teacher": "למורה",
        "sec_student": "לתלמיד",
        "sec_manager": "למנהל/ת",
        "sec_security": "בטיחות ופרטיות",
        "sec_vision": "החזון",

        # Teacher handout
        "t_title": "אנגלית שתלמידים רוצים לעשות",
        "t_subtitle": "Vocaband — כלי דיגיטלי למורי אנגלית, מותאם לכיתה הישראלית",
        "t_problem_label": "הבעיה",
        "t_problem": [
            "תלמידים שוכחים מילים באותו קצב שלומדים אותן",
            "פערים גדולים בכיתה אחת — שיעור אחד לא מתאים לכל התלמידים",
            "אין למורה זמן לתת תרגול אישי ל-30 תלמידים בו זמנית",
            "אחרי שנים של שיבושים — הפערים רק התרחבו",
        ],
        "t_basis_label": "על מה זה בנוי",
        "t_basis": [
            "מדע הלמידה: חזרה מרווחת + גיוון בתרגול",
            "עיצוב משחקי: מוטיבציה פנימית במקום כפיה",
            "מציאות הכיתה: אוטומציה למורה, אנרגיה לתלמיד",
        ],
        "t_get_label": "מה את/ה מקבל/ת",
        "t_get": [
            "בניית מטלות בדקה — לבחור מילים מוכנות או להעלות רשימה משלך",
            "צילום של רשימת מילים מהספר → המערכת מזהה ומתרגמת אוטומטית",
            "בדיקה אוטומטית של תרגילים עם משוב מיידי לתלמיד",
            "דוחות התקדמות לכל תלמיד — מי מתקשה, מי משתעמם, מי מתקדם",
            "Live Challenge — תחרות בזמן אמת על המסך הגדול בכיתה",
            "Quick Play — פעילות כיתתית מהירה עם QR, ללא הרשמה",
            "הקראה ותרגום לעברית ולערבית — בכל מילה, אוטומטית",
        ],
        "t_student_label": "התלמיד יקבל",
        "t_student": [
            "15 משחקים שונים על אותן מילים — בלי שעמום",
            "נקודות, רצפים יומיים, אווטרים, חיות מחמד, תארים",
            "תחרויות כיתתיות חיות — חוויה שמדברים עליה בהפסקה",
        ],
        "t_footer": "ניסיון חינם לכיתה אחת, חודש שלם, ללא התחייבות — אנחנו מקימים את הכיתה איתך",

        # Manager handout
        "m_title": "שליטה, בקרה, בטיחות",
        "m_subtitle": "Vocaband — פלטפורמה ברמת בית ספר, עם נתונים אובייקטיביים",
        "m_problem": [
            "פערי שפה גדולים מאי-פעם בקרב תלמידים בישראל",
            "אין דרך לאתר תלמידים בסיכון לפני שמאוחר",
            "המורים עמוסים — אין נתונים אובייקטיביים על מה קורה בכיתה",
            "הורים ופיקוח דורשים שקיפות — בלי כלי, אין מה להראות",
        ],
        "m_basis": [
            "אותם עקרונות: מדע למידה, גימיפיקציה, אוטומציה למורה",
            "פותח מאפס לבתי ספר בישראל — לא תרגום של מוצר זר",
            "ארכיטקטורת בטיחות מובנית מהיום הראשון — לא טלאי בדיעבד",
        ],
        "m_get": [
            "לוח בקרה אחד לכל הכיתות בבית הספר",
            "נתונים אובייקטיביים: התקדמות, התמדה, אחוזי דיוק",
            "זיהוי מוקדם של תלמידים בסיכון — לפני שמגיעים למבחנים",
            "דוחות מוכנים להורים ולפיקוח — פנים אחת לכל בית הספר",
            "אחסון באירופה (פרנקפורט) — תאימות GDPR + חוק הגנת הפרטיות",
            "ללא פרסומות, ללא מכירת נתונים, ללא רכישות בתוך האפליקציה",
            "נבדק באופן שוטף בבדיקות חדירה (Penetration Tests)",
        ],
        "m_vision_label": "צופה פני עתיד",
        "m_vision": [
            "פלטפורמה אחת — לכל מקצוע שדורש שינון",
            "המנוע מותאם להרחבה לעברית, ערבית, מדעים, היסטוריה ועוד",
            "השקעה בכלי אחד שצומח עם בית הספר",
        ],
        "m_footer": "פיילוט ללא התחייבות — כיתה אחת, חודש שלם, אנחנו עוזרים בהקמה",

        # Deck title slide
        "d_hero": "האפליקציה שתלמידים רוצים לפתוח",
        "d_hero_sub": "פלטפורמה דיגיטלית לרכישת אוצר מילים באנגלית — בנויה לכיתה הישראלית",
        "d_pill": "מצגת לבתי ספר",

        # Slide 2 — problem overview
        "d2_eyebrow": "01 · הבעיה",
        "d2_title": "תלמיד כיתה ט' באוצר מילים של כיתה ו'",
        "d2_cards": [
            ("שכחה מהירה", "תלמידים שוכחים מילים באותו קצב שלומדים אותן. בלי חזרה מגוונת, אוצר מילים מתפוגג תוך שבועות."),
            ("פערים בכיתה", "באותה כיתה — חלק קוראים שוטף, חלק עדיין מתקשים. שיעור אחד לא מתאים ל-30 רמות שונות."),
            ("עומס על המורה", "אין זמן ואין כלים לתת תרגול אישי ולעקוב אחר כל תלמיד בנפרד."),
        ],

        # Slide 3 — problem data
        "d3_eyebrow": "02 · המצב היום",
        "d3_title": "שיטות הלמידה הישנות כבר לא עובדות",
        "d3_points": [
            ("חוברות תרגול", "ילד שמשלים דף מילים פעם בשבוע — שוכח 70% תוך חודש. תרגול לא מגוון לא נשמר בזיכרון."),
            ("שינון בעל-פה", "עובד אצל ילד אחד מתוך עשרה. השאר מתחילים לשנוא את השפה."),
            ("מבחנים בלבד", "מודדים מה התלמיד יודע ביום המבחן — לא מה הוא ידע באמת ובאיזה קצב למד."),
            ("אפליקציות זרות", "מתאימות לתלמיד שכבר יודע אנגלית. לתלמיד שמתחיל מאפס — הן רק מתסכלות."),
        ],

        # Slide 4 — basis overview
        "d4_eyebrow": "03 · היסודות",
        "d4_title": "שלושת העקרונות שעליהם הפלטפורמה בנויה",
        "d4_pillars": [
            ("מדע הלמידה", "חזרה מרווחת (Spaced Repetition) + גיוון בתרגול. אותן מילים, 15 דרכים שונות לתרגל אותן."),
            ("עיצוב משחקי", "מוטיבציה פנימית במקום כפיה. נקודות, רצפים, תארים, אווטרים — מנגנונים שמחזיקים תלמיד שעות."),
            ("מציאות הכיתה", "אוטומציה למורה: בדיקה, ציון, משוב. דוחות לכל תלמיד. בלי להוסיף עוד שעה של עבודה."),
        ],

        # Slide 5 — basis: repetition
        "d5_eyebrow": "04 · עיקרון ראשון",
        "d5_title": "אותה מילה — 15 דרכים שונות לתרגל",
        "d5_subhead": "למה זה עובד",
        "d5_bullets": [
            "התלמיד פוגש את המילה כקריאה, ככתיבה, כשמיעה ובהקשר משפט",
            "כל מצב משחק מפעיל אזור שונה במוח — הזיכרון מתחזק",
            "המעבר בין משחקים שובר שעמום ושומר על ריכוז",
            "התלמיד לא מרגיש שהוא עושה את אותו דבר שוב",
            "המורה לא צריך/ה להמציא תרגילים חדשים",
        ],

        # Slide 6 — basis: motivation
        "d6_eyebrow": "05 · עיקרון שני",
        "d6_title": "מוטיבציה פנימית — כי כך התלמיד חוזר מחר",
        "d6_bullets": [
            "כל תשובה נכונה משלמת נקודות XP",
            "רצף יומי שמתגמל התמדה, לא רק הצלחה",
            "חנות אווטרים, תארים ופריימים — נרכשים בנקודות בלבד",
            "חיית מחמד שמתפתחת לאורך זמן",
            "תחרויות חיות מול חברים לכיתה",
            "אין פרסומות, אין רכישות בכסף אמיתי, אין לחץ",
        ],

        # Slide 7 — basis: automation
        "d7_eyebrow": "06 · עיקרון שלישי",
        "d7_title": "אוטומציה למורה — בלי להוסיף שעות עבודה",
        "d7_bullets": [
            "מטלות נבדקות אוטומטית, התלמיד מקבל משוב מיידי",
            "המערכת מציינת את התלמיד ושומרת היסטוריה מלאה",
            "דוחות התקדמות נוצרים אוטומטית — לפי תלמיד, לפי כיתה, לפי תקופה",
            "המורה רואה/ה מי בקושי, מי משתעמם, מי מתקדם — בלי לחפש",
            "תזכורות אוטומטיות לתלמידים שלא נכנסו השבוע",
        ],

        # Slide 8 — teachers
        "d8_eyebrow": "07 · למורה",
        "d8_title": "פחות בדיקות, יותר הוראה",
        "d8_bullets": [
            "בניית מטלה בדקה — מילים מוכנות או רשימה משלך",
            "צילום של עמוד מהספר — זיהוי אוטומטי + תרגום",
            "בדיקה אוטומטית של תרגילים",
            "דוחות התקדמות פר תלמיד",
            "Live Challenge — תחרות בזמן אמת על המסך הגדול",
            "Quick Play — פעילות מהירה עם QR, בלי הרשמה",
            "הקראה ותרגום לעברית/ערבית בכל מילה",
            "כניסת תלמידים עם קוד כיתה — בלי תהליך הרשמה מסובך",
        ],

        # Slide 9 — students
        "d9_eyebrow": "08 · לתלמיד",
        "d9_title": "תרגול שמרגיש כמו משחק, לא כמו שיעורי בית",
        "d9_bullets": [
            "15 מצבי משחק שונים על אותן מילים",
            "נקודות XP, רצפים יומיים, רמות, תארים",
            "חנות פנים-משחק עם אווטרים, פריימים, ערכות צבע",
            "חיית מחמד שמתפתחת לאורך זמן",
            "מנגנוני שימור: תיבת יומית, אתגר שבועי, בונוס חזרה",
            "תחרות חיה מול חברים לכיתה",
            "תרגום מיידי לעברית או ערבית כשנתקעים",
            "הקראה לכל מילה — בהגייה נכונה",
        ],

        # Slide 10 — managers
        "d10_eyebrow": "09 · למנהל/ת",
        "d10_title": "שליטה, בקרה, נתונים — לא תחושות בטן",
        "d10_bullets": [
            "לוח בקרה אחד לכל הכיתות בבית הספר",
            "נתונים אובייקטיביים על כל תלמיד וכל כיתה",
            "זיהוי מוקדם של תלמידים בסיכון",
            "דוחות מוכנים להורים ולפיקוח",
            "השוואה בין כיתות ולאורך זמן",
            "פלטפורמה אחת, ללא תוכנה להתקין",
            "ללא פרסומות, ללא רכישות בתוך האפליקציה",
            "ניתן להרחבה לכל בית הספר ללא תוספת תשתית",
        ],

        # Slide 11 — security
        "d11_eyebrow": "10 · בטיחות ופרטיות",
        "d11_title": "תאימות מלאה לחוק הישראלי ולתקן האירופי",
        "d11_cards": [
            ("אחסון באירופה", "כל הנתונים מאוחסנים בפרנקפורט. תאימות מלאה ל-GDPR."),
            ("חוק הגנת הפרטיות", "עומדים בחוק הגנת הפרטיות הישראלי ובתיקון 13."),
            ("ללא פרסומות", "אין פרסומות, אין צד שלישי, אין מכירת נתונים."),
            ("בדיקות חדירה", "המערכת נבדקת באופן שוטף בבדיקות חדירה חיצוניות."),
        ],

        # Slide 12 — vision
        "d12_eyebrow": "11 · החזון",
        "d12_title": "אנגלית היום — כל מקצוע שדורש שינון, מחר",
        "d12_desc": "המנוע שבנינו לאוצר מילים באנגלית מתאים לכל תחום שדורש שינון והפנמה.",
        "d12_pills": [
            ("עברית כשפה שנייה", INDIGO),
            ("ערבית", VIOLET),
            ("אוצר מילים במדעים", EMERALD),
            ("מושגים בהיסטוריה ובתנ\"ך", AMBER),
            ("שפה שלישית — צרפתית, ספרדית, רוסית", FUCHSIA),
        ],

        # Slide 13 — how to start
        "d13_hero": "איך מתחילים",
        "d13_steps": [
            ("01", "פגישת היכרות קצרה", "אנחנו לומדים את בית הספר ואת הצרכים שלכם"),
            ("02", "בוחרים כיתה אחת לפיילוט", "אתם בוחרים מורה, אנחנו מקימים את הכיתה"),
            ("03", "חודש שלם של שימוש", "ללא התחייבות, ללא תשלום, ליווי צמוד"),
            ("04", "מחליטים יחד", "אם זה עובד — מרחיבים. אם לא — לא ממשיכים."),
        ],

        # Slide 14 — thanks
        "d14_thanks": "תודה.",
        "d14_sub": "נשמח להראות לכם את המערכת בפגישה קצרה.",
    },

    "ar": {
        "tag": "AR",
        "font": "Noto Sans Arabic",

        # Page chrome / section labels
        "sec_problem": "المشكلة",
        "sec_basis": "الأسس",
        "sec_teacher": "للمعلم/ة",
        "sec_student": "للطالب",
        "sec_manager": "للمدير/ة",
        "sec_security": "أمان وخصوصية",
        "sec_vision": "الرؤية",

        # Teacher handout
        "t_title": "إنجليزية يحب الطلاب أن يتعلموها",
        "t_subtitle": "Vocaband — أداة رقمية لمعلمي اللغة الإنجليزية، مُصمَّمة للصف الإسرائيلي",
        "t_problem_label": "المشكلة",
        "t_problem": [
            "الطلاب ينسون الكلمات بنفس سرعة تعلُّمها",
            "فجوات كبيرة داخل الصف الواحد — درس واحد لا يناسب جميع الطلاب",
            "لا يجد المعلم/ة الوقت لتقديم تدريب فردي لـ 30 طالبًا في آنٍ واحد",
            "بعد سنوات من الاضطرابات — اتسعت الفجوات أكثر",
        ],
        "t_basis_label": "الأسس التي بُني عليها",
        "t_basis": [
            "علم التعلّم: التكرار المتباعد + تنويع التدريب",
            "تصميم لعبة: دافعية داخلية بدلًا من الإلزام",
            "واقع الصف: أتمتة للمعلم/ة، طاقة للطالب",
        ],
        "t_get_label": "ما الذي تحصل/ين عليه",
        "t_get": [
            "بناء مهام في دقيقة — اختاري كلمات جاهزة أو ارفعي قائمتك",
            "صورة لقائمة كلمات من الكتاب → النظام يتعرف عليها ويترجمها تلقائيًا",
            "تصحيح آلي للتمارين مع تغذية راجعة فورية للطالب",
            "تقارير تقدُّم لكل طالب — من يعاني، من يمل، من يتقدم",
            "Live Challenge — منافسة بالوقت الحقيقي على الشاشة الكبيرة في الصف",
            "Quick Play — نشاط صفّي سريع برمز QR، دون تسجيل",
            "نطق وترجمة للعبرية والعربية — لكل كلمة، تلقائيًا",
        ],
        "t_student_label": "ما يحصل عليه الطالب",
        "t_student": [
            "15 لعبة مختلفة على الكلمات نفسها — دون ملل",
            "نقاط، سلاسل يومية، أفاتارات، حيوانات أليفة، ألقاب",
            "منافسات صفّية حية — تجربة يتحدثون عنها في الاستراحة",
        ],
        "t_footer": "تجربة مجانية لصف واحد، شهر كامل، دون التزام — نُعدّ الصف معك",

        # Manager handout
        "m_title": "تحكُّم، رقابة، أمان",
        "m_subtitle": "Vocaband — منصة بمستوى المدرسة، ببيانات موضوعية",
        "m_problem": [
            "فجوات لغوية كبيرة كما لم تكن من قبل بين الطلاب في إسرائيل",
            "لا توجد طريقة لرصد الطلاب المعرَّضين للخطر قبل فوات الأوان",
            "المعلمون مرهقون — لا توجد بيانات موضوعية عمّا يجري في الصف",
            "الأهل والإشراف يطالبون بالشفافية — بلا أداة، لا شيء يُعرض",
        ],
        "m_basis": [
            "نفس المبادئ: علم التعلّم، التلعيب، الأتمتة للمعلم/ة",
            "طُوِّر من الصفر لمدارس إسرائيل — وليس ترجمة لمنتج أجنبي",
            "بنية أمان مدمجة منذ اليوم الأول — وليست رقعة لاحقة",
        ],
        "m_get": [
            "لوحة تحكم واحدة لجميع صفوف المدرسة",
            "بيانات موضوعية: تقدُّم، مثابرة، نسب دقة",
            "رصد مبكر للطلاب المعرَّضين للخطر — قبل الوصول إلى الامتحانات",
            "تقارير جاهزة للأهل وللإشراف — صورة واحدة لكل المدرسة",
            "تخزين في أوروبا (فرانكفورت) — توافق GDPR + قانون حماية الخصوصية",
            "بدون إعلانات، بدون بيع بيانات، بدون مشتريات داخل التطبيق",
            "تُختبر بانتظام في اختبارات الاختراق (Penetration Tests)",
        ],
        "m_vision_label": "نظرة مستقبلية",
        "m_vision": [
            "منصة واحدة — لكل مادة تتطلب الحفظ",
            "المحرك مُهيَّأ للتوسع إلى العبرية، العربية، العلوم، التاريخ والمزيد",
            "استثمار في أداة واحدة تنمو مع المدرسة",
        ],
        "m_footer": "تجربة دون التزام — صف واحد، شهر كامل، نساعد في الإعداد",

        # Deck title slide
        "d_hero": "التطبيق الذي يحب الطلاب فتحه",
        "d_hero_sub": "منصة رقمية لاكتساب المفردات بالإنجليزية — مصممة للصف الإسرائيلي",
        "d_pill": "عرض تقديمي للمدارس",

        # Slide 2 — problem overview
        "d2_eyebrow": "01 · المشكلة",
        "d2_title": "طالب في الصف التاسع بمفردات الصف السادس",
        "d2_cards": [
            ("نسيان سريع", "ينسى الطلاب الكلمات بنفس سرعة تعلُّمها. دون تكرار متنوع، تتلاشى المفردات خلال أسابيع."),
            ("فجوات في الصف", "في نفس الصف — بعضهم يقرأ بطلاقة، وبعضهم لا يزال يكافح. درس واحد لا يناسب 30 مستوى مختلفًا."),
            ("عبء على المعلم/ة", "لا وقت ولا أدوات لتقديم تدريب فردي ومتابعة كل طالب على حدة."),
        ],

        # Slide 3 — problem data
        "d3_eyebrow": "02 · الواقع اليوم",
        "d3_title": "أساليب التعلّم القديمة لم تعد ناجعة",
        "d3_points": [
            ("كرّاسات التدريب", "طفل يكمل ورقة مفردات مرة في الأسبوع — ينسى 70% خلال شهر. التدريب غير المتنوع لا يبقى في الذاكرة."),
            ("الحفظ غيبًا", "ينجح مع طفل من عشرة. الباقون يبدؤون بكره اللغة."),
            ("الامتحانات فقط", "تقيس ما يعرفه الطالب يوم الامتحان — لا ما عرفه فعلًا وبأي وتيرة تعلَّم."),
            ("تطبيقات أجنبية", "تناسب الطالب الذي يعرف الإنجليزية مسبقًا. أمّا المبتدئ من الصفر — فهي تُشعره بالإحباط."),
        ],

        # Slide 4 — basis overview
        "d4_eyebrow": "03 · الأسس",
        "d4_title": "المبادئ الثلاثة التي بُنيت عليها المنصة",
        "d4_pillars": [
            ("علم التعلّم", "التكرار المتباعد (Spaced Repetition) + تنويع التدريب. نفس الكلمات، 15 طريقة مختلفة للتدرّب عليها."),
            ("تصميم لعبة", "دافعية داخلية بدلًا من الإلزام. نقاط، سلاسل، ألقاب، أفاتارات — آليات تُبقي الطالب ساعات."),
            ("واقع الصف", "أتمتة للمعلم/ة: تصحيح، علامة، تغذية راجعة. تقارير لكل طالب. دون إضافة ساعة عمل أخرى."),
        ],

        # Slide 5 — basis: repetition
        "d5_eyebrow": "04 · المبدأ الأول",
        "d5_title": "نفس الكلمة — 15 طريقة مختلفة للتدريب",
        "d5_subhead": "لماذا ينجح هذا",
        "d5_bullets": [
            "يلتقي الطالب بالكلمة قراءةً وكتابةً وسماعًا وفي سياق جملة",
            "كل وضع لعب يُشغّل منطقة مختلفة في الدماغ — تتقوّى الذاكرة",
            "التنقّل بين الألعاب يكسر الملل ويحافظ على التركيز",
            "لا يشعر الطالب أنه يكرر الشيء نفسه",
            "لا يحتاج المعلم/ة إلى ابتكار تمارين جديدة",
        ],

        # Slide 6 — basis: motivation
        "d6_eyebrow": "05 · المبدأ الثاني",
        "d6_title": "دافعية داخلية — حتى يعود الطالب غدًا",
        "d6_bullets": [
            "كل إجابة صحيحة تمنح نقاط XP",
            "سلسلة يومية تكافئ المثابرة، لا النجاح فحسب",
            "متجر أفاتارات وألقاب وإطارات — تُشترى بالنقاط فقط",
            "حيوان أليف يتطور مع مرور الوقت",
            "منافسات حية مع زملاء الصف",
            "بلا إعلانات، بلا مشتريات بمال حقيقي، بلا ضغط",
        ],

        # Slide 7 — basis: automation
        "d7_eyebrow": "06 · المبدأ الثالث",
        "d7_title": "أتمتة للمعلم/ة — دون إضافة ساعات عمل",
        "d7_bullets": [
            "تُصحَّح المهام تلقائيًا، ويحصل الطالب على تغذية راجعة فورية",
            "النظام يضع العلامة ويحفظ سجلًا كاملًا",
            "تقارير التقدم تُنشأ تلقائيًا — لكل طالب، لكل صف، لكل فترة",
            "يرى المعلم/ة من يعاني، من يمل، من يتقدم — دون بحث",
            "تذكيرات تلقائية للطلاب الذين لم يدخلوا هذا الأسبوع",
        ],

        # Slide 8 — teachers
        "d8_eyebrow": "07 · للمعلم/ة",
        "d8_title": "تصحيح أقل، تدريس أكثر",
        "d8_bullets": [
            "بناء مهمة في دقيقة — كلمات جاهزة أو قائمتك",
            "صورة لصفحة من الكتاب — تعرّف تلقائي + ترجمة",
            "تصحيح آلي للتمارين",
            "تقارير تقدم لكل طالب",
            "Live Challenge — منافسة بالوقت الحقيقي على الشاشة الكبيرة",
            "Quick Play — نشاط سريع برمز QR، دون تسجيل",
            "نطق وترجمة للعبرية/العربية لكل كلمة",
            "دخول الطلاب برمز الصف — دون عملية تسجيل معقدة",
        ],

        # Slide 9 — students
        "d9_eyebrow": "08 · للطالب",
        "d9_title": "تدريب يُحَسّ كأنه لعبة، لا كواجب بيتي",
        "d9_bullets": [
            "15 وضع لعب مختلف على الكلمات نفسها",
            "نقاط XP، سلاسل يومية، مستويات، ألقاب",
            "متجر داخل اللعبة بأفاتارات وإطارات وحزم ألوان",
            "حيوان أليف يتطور مع مرور الوقت",
            "آليات استبقاء: صندوق يومي، تحدٍّ أسبوعي، مكافأة عودة",
            "منافسة حية مع زملاء الصف",
            "ترجمة فورية للعبرية أو العربية عند التعثُّر",
            "نطق لكل كلمة — بنطق صحيح",
        ],

        # Slide 10 — managers
        "d10_eyebrow": "09 · للمدير/ة",
        "d10_title": "تحكم، رقابة، بيانات — لا أحاسيس",
        "d10_bullets": [
            "لوحة تحكم واحدة لجميع صفوف المدرسة",
            "بيانات موضوعية عن كل طالب وكل صف",
            "رصد مبكر للطلاب المعرَّضين للخطر",
            "تقارير جاهزة للأهل وللإشراف",
            "مقارنة بين الصفوف وعلى مرّ الزمن",
            "منصة واحدة، دون برامج للتثبيت",
            "بلا إعلانات، بلا مشتريات داخل التطبيق",
            "قابلة للتوسع لكل المدرسة دون بنية تحتية إضافية",
        ],

        # Slide 11 — security
        "d11_eyebrow": "10 · أمان وخصوصية",
        "d11_title": "توافق كامل مع القانون الإسرائيلي والمعيار الأوروبي",
        "d11_cards": [
            ("تخزين في أوروبا", "تُخزَّن جميع البيانات في فرانكفورت. توافق كامل مع GDPR."),
            ("قانون حماية الخصوصية", "نتقيد بقانون حماية الخصوصية الإسرائيلي وبالتعديل 13."),
            ("بلا إعلانات", "لا إعلانات، لا طرف ثالث، لا بيع للبيانات."),
            ("اختبارات اختراق", "يُختبر النظام بانتظام في اختبارات اختراق خارجية."),
        ],

        # Slide 12 — vision
        "d12_eyebrow": "11 · الرؤية",
        "d12_title": "الإنجليزية اليوم — أي مادة تتطلب الحفظ، غدًا",
        "d12_desc": "المحرك الذي بنيناه لمفردات الإنجليزية يناسب أي مجال يتطلب الحفظ والاستيعاب.",
        "d12_pills": [
            ("العبرية كلغة ثانية", INDIGO),
            ("العربية", VIOLET),
            ("مفردات في العلوم", EMERALD),
            ("مفاهيم في التاريخ والتناخ", AMBER),
            ("لغة ثالثة — فرنسية، إسبانية، روسية", FUCHSIA),
        ],

        # Slide 13 — how to start
        "d13_hero": "كيف نبدأ",
        "d13_steps": [
            ("01", "لقاء تعارف قصير", "نتعرّف على المدرسة واحتياجاتكم"),
            ("02", "نختار صفًا واحدًا للتجربة", "تختارون المعلم/ة، ونحن نُعدّ الصف"),
            ("03", "شهر كامل من الاستخدام", "دون التزام، دون دفع، بمتابعة قريبة"),
            ("04", "نقرر معًا", "إن نجحت — نوسّع. وإلا — لا نُكمل."),
        ],

        # Slide 14 — thanks
        "d14_thanks": "شكرًا.",
        "d14_sub": "يسعدنا أن نُريكم النظام في لقاء قصير.",
    },
}


# ---------- Low-level helpers ----------

def set_rtl(paragraph) -> None:
    """Mark a paragraph as RTL so punctuation and digits land correctly."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.RIGHT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if align == PP_ALIGN.RIGHT:
            set_rtl(p)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullet_list(slide, x, y, w, h, items, *, font, size=12, color=INK, marker="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.line_spacing = 1.4
        p.space_after = Pt(2)
        set_rtl(p)
        # RTL paragraph: marker first in source -> visual right edge.
        run = p.add_run()
        run.text = f"{marker}  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_image(slide, x, y, w, h, path):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


# ---------- A4 portrait one-pager helpers ----------

def make_a4_portrait_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    return prs


def handout_header(slide, title, subtitle, accent_color, font):
    add_rect(slide, 0, 0, Cm(21.0), Cm(4.2), accent_color)
    add_text(
        slide, Cm(0.8), Cm(0.6), Cm(19.4), Cm(0.9),
        "VOCABAND",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        font="DejaVu Sans",
    )
    add_text(
        slide, Cm(0.8), Cm(1.4), Cm(19.4), Cm(1.6),
        title,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
        font=font,
    )
    add_text(
        slide, Cm(0.8), Cm(3.0), Cm(19.4), Cm(1.0),
        subtitle,
        size=14, color=WHITE, align=PP_ALIGN.RIGHT,
        font=font,
    )


def handout_section(slide, y_cm, label, label_color, body_lines, font):
    pill = add_rounded(
        slide,
        Cm(21.0 - 0.8 - 5.0), Cm(y_cm),
        Cm(5.0), Cm(0.85),
        label_color,
    )
    tf = pill.text_frame
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_rtl(p)
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_bullet_list(
        slide,
        Cm(0.8), Cm(y_cm + 1.05), Cm(19.4), Cm(6.0),
        body_lines,
        size=12,
        color=INK,
        font=font,
    )


def handout_footer(slide, text, color, font):
    add_rect(slide, 0, Cm(27.5), Cm(21.0), Cm(2.2), color)
    add_text(
        slide, Cm(0.8), Cm(27.85), Cm(19.4), Cm(1.5),
        text,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        font=font,
    )


# ---------- Handouts ----------

def build_teacher_handout(s):
    font = s["font"]
    prs = make_a4_portrait_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    handout_header(slide, s["t_title"], s["t_subtitle"], INDIGO, font)
    handout_section(slide, 4.8,  s["t_problem_label"], ROSE,    s["t_problem"], font)
    handout_section(slide, 10.6, s["t_basis_label"],   VIOLET,  s["t_basis"], font)
    handout_section(slide, 15.1, s["t_get_label"],     INDIGO,  s["t_get"], font)
    handout_section(slide, 23.2, s["t_student_label"], AMBER,   s["t_student"], font)
    handout_footer(slide, s["t_footer"], INK, font)

    out = OUT_DIR / f"teacher-handout-{s['tag']}.pptx"
    prs.save(out)
    return out


def build_manager_handout(s):
    font = s["font"]
    prs = make_a4_portrait_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    handout_header(slide, s["m_title"], s["m_subtitle"], VIOLET, font)
    handout_section(slide, 4.8,  s["t_problem_label"], ROSE,    s["m_problem"], font)
    handout_section(slide, 10.6, s["t_basis_label"],   VIOLET,  s["m_basis"], font)
    handout_section(slide, 15.1, s["t_get_label"],     INDIGO,  s["m_get"], font)
    handout_section(slide, 23.2, s["m_vision_label"],  EMERALD, s["m_vision"], font)
    handout_footer(slide, s["m_footer"], INK, font)

    out = OUT_DIR / f"manager-handout-{s['tag']}.pptx"
    prs.save(out)
    return out


# ---------- Full presentation ----------

def make_widescreen_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def deck_background(slide, prs):
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)


def deck_side_strip(slide, prs, color):
    add_rect(slide, prs.slide_width - Cm(0.5), 0, Cm(0.5), prs.slide_height, color)


def deck_page_chrome(slide, prs, page_num, total, section_label, color, font):
    add_rect(slide, 0, prs.slide_height - Cm(0.6), prs.slide_width, Cm(0.6), BG_SOFT)
    add_text(
        slide, Cm(0.8), prs.slide_height - Cm(0.55), Cm(20.0), Cm(0.5),
        section_label,
        size=10, bold=True, color=color, align=PP_ALIGN.RIGHT, font=font,
    )
    add_text(
        slide, Cm(24.0), prs.slide_height - Cm(0.55), Cm(8.5), Cm(0.5),
        f"{page_num} / {total}",
        size=10, color=MUTED, align=PP_ALIGN.LEFT, font="DejaVu Sans",
    )
    add_text(
        slide, Cm(0.8), Cm(0.4), Cm(8.0), Cm(0.6),
        "VOCABAND",
        size=11, bold=True, color=color, align=PP_ALIGN.LEFT, font="DejaVu Sans",
    )


def deck_title(slide, prs, eyebrow, title, color, font):
    add_text(
        slide, Cm(0.8), Cm(1.4), Cm(32.0), Cm(0.8),
        eyebrow,
        size=12, bold=True, color=color, align=PP_ALIGN.RIGHT, font=font,
    )
    add_text(
        slide, Cm(0.8), Cm(2.1), Cm(32.0), Cm(2.0),
        title,
        size=34, bold=True, color=INK, align=PP_ALIGN.RIGHT, font=font,
    )


def slide_title_page(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, INDIGO)
    add_rect(sl, Cm(11.3), 0, Cm(11.3), prs.slide_height, VIOLET)
    add_rect(sl, Cm(22.6), 0, Cm(11.3), prs.slide_height, FUCHSIA)

    add_text(
        sl, Cm(0.8), Cm(0.8), Cm(32.0), Cm(0.8),
        "VOCABAND",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="DejaVu Sans",
    )
    add_text(
        sl, Cm(1.5), Cm(6.5), Cm(31.0), Cm(2.2),
        s["d_hero"],
        size=44, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
    )
    add_text(
        sl, Cm(1.5), Cm(9.0), Cm(31.0), Cm(1.5),
        s["d_hero_sub"],
        size=18, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
    )
    add_rounded(sl, Cm(1.5), Cm(14.5), Cm(8.0), Cm(1.2), WHITE)
    add_text(
        sl, Cm(1.5), Cm(14.5), Cm(8.0), Cm(1.2),
        s["d_pill"],
        size=14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=font,
    )


def slide_problem_overview(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, ROSE)
    deck_page_chrome(sl, prs, 2, 14, s["sec_problem"], ROSE, font)
    deck_title(sl, prs, s["d2_eyebrow"], s["d2_title"], ROSE, font)

    colors = [ROSE, AMBER, VIOLET]
    card_w = Cm(10.0)
    card_h = Cm(9.0)
    gap = Cm(0.8)
    start_x = (prs.slide_width - (card_w * 3 + gap * 2)) / 2

    for i, (title, body) in enumerate(s["d2_cards"]):
        x = start_x + (card_w + gap) * i
        y = Cm(5.5)
        add_rounded(sl, x, y, card_w, card_h, BG_SOFT)
        add_rounded(sl, x, y, card_w, Cm(0.4), colors[i])
        add_text(
            sl, x + Cm(0.5), y + Cm(0.9), card_w - Cm(1.0), Cm(1.2),
            title, size=20, bold=True, color=INK, align=PP_ALIGN.RIGHT, font=font,
        )
        add_text(
            sl, x + Cm(0.5), y + Cm(2.5), card_w - Cm(1.0), Cm(6.0),
            body, size=14, color=SLATE, align=PP_ALIGN.RIGHT,
            line_spacing=1.35, font=font,
        )


def slide_problem_data(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, ROSE)
    deck_page_chrome(sl, prs, 3, 14, s["sec_problem"], ROSE, font)
    deck_title(sl, prs, s["d3_eyebrow"], s["d3_title"], ROSE, font)

    y = Cm(5.5)
    for title, body in s["d3_points"]:
        add_rect(sl, Cm(0.8), y, Cm(32.0), Cm(0.05), BG_SOFT)
        add_text(
            sl, Cm(0.8), y + Cm(0.25), Cm(32.0), Cm(0.8),
            title, size=18, bold=True, color=ROSE, align=PP_ALIGN.RIGHT, font=font,
        )
        add_text(
            sl, Cm(0.8), y + Cm(1.1), Cm(32.0), Cm(1.4),
            body, size=14, color=SLATE, align=PP_ALIGN.RIGHT, font=font,
        )
        y += Cm(2.8)


def slide_basis_overview(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, VIOLET)
    deck_page_chrome(sl, prs, 4, 14, s["sec_basis"], VIOLET, font)
    deck_title(sl, prs, s["d4_eyebrow"], s["d4_title"], VIOLET, font)

    colors = [INDIGO, FUCHSIA, EMERALD]
    card_w = Cm(10.0)
    card_h = Cm(9.5)
    gap = Cm(0.8)
    start_x = (prs.slide_width - (card_w * 3 + gap * 2)) / 2

    for i, (title, body) in enumerate(s["d4_pillars"]):
        x = start_x + (card_w + gap) * i
        y = Cm(5.3)
        add_rounded(sl, x, y, card_w, card_h, colors[i])
        add_rounded(sl, x + card_w - Cm(2.2), y + Cm(0.6), Cm(1.6), Cm(1.6), WHITE)
        add_text(
            sl, x + card_w - Cm(2.2), y + Cm(0.6), Cm(1.6), Cm(1.6),
            f"0{i+1}", size=24, bold=True, color=colors[i], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font="DejaVu Sans",
        )
        add_text(
            sl, x + Cm(0.6), y + Cm(2.8), card_w - Cm(1.2), Cm(1.2),
            title, size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
        )
        add_text(
            sl, x + Cm(0.6), y + Cm(4.3), card_w - Cm(1.2), Cm(5.0),
            body, size=14, color=WHITE, align=PP_ALIGN.RIGHT, line_spacing=1.4, font=font,
        )


def slide_basis_repetition(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, VIOLET)
    deck_page_chrome(sl, prs, 5, 14, s["sec_basis"], VIOLET, font)
    deck_title(sl, prs, s["d5_eyebrow"], s["d5_title"], VIOLET, font)

    add_text(
        sl, Cm(0.8), Cm(5.5), Cm(18.0), Cm(1.0),
        s["d5_subhead"],
        size=18, bold=True, color=VIOLET, align=PP_ALIGN.RIGHT, font=font,
    )
    add_bullet_list(
        sl, Cm(0.8), Cm(6.7), Cm(18.0), Cm(10.0),
        s["d5_bullets"], size=14, font=font,
    )

    img = SHOTS / "screen2.png"
    if img.exists():
        add_image(sl, Cm(21.5), Cm(5.0), Cm(8.5), Cm(13.0), img)


def slide_basis_motivation(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, VIOLET)
    deck_page_chrome(sl, prs, 6, 14, s["sec_basis"], VIOLET, font)
    deck_title(sl, prs, s["d6_eyebrow"], s["d6_title"], FUCHSIA, font)

    add_bullet_list(
        sl, Cm(0.8), Cm(6.0), Cm(18.0), Cm(11.0),
        s["d6_bullets"], size=14, font=font,
    )
    img = SHOTS / "screen.png"
    if img.exists():
        add_image(sl, Cm(21.5), Cm(5.0), Cm(8.5), Cm(13.0), img)


def slide_basis_automation(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, VIOLET)
    deck_page_chrome(sl, prs, 7, 14, s["sec_basis"], VIOLET, font)
    deck_title(sl, prs, s["d7_eyebrow"], s["d7_title"], EMERALD, font)

    add_bullet_list(
        sl, Cm(0.8), Cm(6.0), Cm(32.0), Cm(11.0),
        s["d7_bullets"], size=15, font=font,
    )


def slide_for_teachers(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, INDIGO)
    deck_page_chrome(sl, prs, 8, 14, s["sec_teacher"], INDIGO, font)
    deck_title(sl, prs, s["d8_eyebrow"], s["d8_title"], INDIGO, font)

    add_bullet_list(
        sl, Cm(0.8), Cm(5.5), Cm(18.0), Cm(12.0),
        s["d8_bullets"], size=14, font=font,
    )
    img = SHOTS / "screen1.png"
    if img.exists():
        add_image(sl, Cm(21.5), Cm(5.0), Cm(8.5), Cm(13.0), img)


def slide_for_students(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, AMBER)
    deck_page_chrome(sl, prs, 9, 14, s["sec_student"], AMBER, font)
    deck_title(sl, prs, s["d9_eyebrow"], s["d9_title"], AMBER, font)

    add_bullet_list(
        sl, Cm(0.8), Cm(5.5), Cm(18.0), Cm(12.0),
        s["d9_bullets"], size=14, font=font,
    )
    img = SHOTS / "screen.png"
    if img.exists():
        add_image(sl, Cm(21.5), Cm(5.0), Cm(8.5), Cm(13.0), img)


def slide_for_managers(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, FUCHSIA)
    deck_page_chrome(sl, prs, 10, 14, s["sec_manager"], FUCHSIA, font)
    deck_title(sl, prs, s["d10_eyebrow"], s["d10_title"], FUCHSIA, font)

    add_bullet_list(
        sl, Cm(0.8), Cm(5.5), Cm(18.0), Cm(12.0),
        s["d10_bullets"], size=14, font=font,
    )
    img = SHOTS / "scree3n.png"
    if img.exists():
        add_image(sl, Cm(21.5), Cm(5.0), Cm(8.5), Cm(13.0), img)


def slide_security(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, EMERALD)
    deck_page_chrome(sl, prs, 11, 14, s["sec_security"], EMERALD, font)
    deck_title(sl, prs, s["d11_eyebrow"], s["d11_title"], EMERALD, font)

    card_w = Cm(15.5)
    card_h = Cm(5.5)
    gap = Cm(0.8)
    start_x = (prs.slide_width - (card_w * 2 + gap)) / 2

    for i, (title, body) in enumerate(s["d11_cards"]):
        col = i % 2
        row = i // 2
        x = start_x + (card_w + gap) * col
        y = Cm(5.5) + (card_h + gap) * row
        add_rounded(sl, x, y, card_w, card_h, BG_SOFT)
        add_rounded(sl, x, y, Cm(0.4), card_h, EMERALD)
        add_text(
            sl, x + Cm(0.7), y + Cm(0.7), card_w - Cm(1.4), Cm(1.2),
            title, size=18, bold=True, color=INK, align=PP_ALIGN.RIGHT, font=font,
        )
        add_text(
            sl, x + Cm(0.7), y + Cm(2.2), card_w - Cm(1.4), Cm(3.0),
            body, size=14, color=SLATE, align=PP_ALIGN.RIGHT,
            line_spacing=1.4, font=font,
        )


def slide_vision(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    deck_background(sl, prs)
    deck_side_strip(sl, prs, FUCHSIA)
    deck_page_chrome(sl, prs, 12, 14, s["sec_vision"], FUCHSIA, font)
    deck_title(sl, prs, s["d12_eyebrow"], s["d12_title"], FUCHSIA, font)

    add_text(
        sl, Cm(0.8), Cm(5.5), Cm(32.0), Cm(2.0),
        s["d12_desc"],
        size=18, color=SLATE, align=PP_ALIGN.RIGHT, line_spacing=1.4, font=font,
    )

    y = Cm(9.0)
    for label, color in s["d12_pills"]:
        add_rounded(sl, Cm(0.8), y, Cm(32.0), Cm(1.3), color)
        add_text(
            sl, Cm(0.8), y, Cm(32.0), Cm(1.3),
            label, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font=font,
        )
        y += Cm(1.55)


def slide_start(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, INDIGO)
    add_rect(sl, Cm(11.3), 0, Cm(22.6), prs.slide_height, VIOLET)
    add_rect(sl, Cm(22.6), 0, Cm(11.3), prs.slide_height, FUCHSIA)

    add_text(
        sl, Cm(0.8), Cm(0.8), Cm(32.0), Cm(0.8),
        "VOCABAND",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="DejaVu Sans",
    )
    add_text(
        sl, Cm(1.5), Cm(5.5), Cm(31.0), Cm(2.2),
        s["d13_hero"],
        size=42, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
    )

    y = Cm(8.5)
    for num, title, body in s["d13_steps"]:
        add_rounded(sl, Cm(1.5), y, Cm(2.5), Cm(2.0), WHITE)
        add_text(
            sl, Cm(1.5), y, Cm(2.5), Cm(2.0),
            num, size=26, bold=True, color=INDIGO, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font="DejaVu Sans",
        )
        add_text(
            sl, Cm(4.3), y, Cm(28.0), Cm(0.9),
            title, size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
        )
        add_text(
            sl, Cm(4.3), y + Cm(1.0), Cm(28.0), Cm(1.0),
            body, size=13, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
        )
        y += Cm(2.3)


def slide_thanks(prs, s):
    font = s["font"]
    sl = slide_blank(prs)
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, INK)

    add_text(
        sl, Cm(0.8), Cm(0.8), Cm(32.0), Cm(0.8),
        "VOCABAND",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="DejaVu Sans",
    )
    add_text(
        sl, Cm(1.5), Cm(7.0), Cm(31.0), Cm(3.0),
        s["d14_thanks"],
        size=72, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=font,
    )
    add_text(
        sl, Cm(1.5), Cm(11.5), Cm(31.0), Cm(1.5),
        s["d14_sub"],
        size=18, color=MUTED, align=PP_ALIGN.RIGHT, font=font,
    )
    add_text(
        sl, Cm(1.5), Cm(13.5), Cm(31.0), Cm(1.0),
        "www.vocaband.com",
        size=18, bold=True, color=AMBER, align=PP_ALIGN.RIGHT, font="DejaVu Sans",
    )


def build_presentation(s):
    prs = make_widescreen_prs()
    slide_title_page(prs, s)
    slide_problem_overview(prs, s)
    slide_problem_data(prs, s)
    slide_basis_overview(prs, s)
    slide_basis_repetition(prs, s)
    slide_basis_motivation(prs, s)
    slide_basis_automation(prs, s)
    slide_for_teachers(prs, s)
    slide_for_students(prs, s)
    slide_for_managers(prs, s)
    slide_security(prs, s)
    slide_vision(prs, s)
    slide_start(prs, s)
    slide_thanks(prs, s)

    out = OUT_DIR / f"vocaband-presentation-{s['tag']}.pptx"
    prs.save(out)
    return out


# ---------- PDF conversion ----------

def to_pdf(pptx_path: Path) -> Path:
    pdf_path = pptx_path.with_suffix(".pdf")
    if pdf_path.exists():
        pdf_path.unlink()
    result = subprocess.run(
        [
            "libreoffice",
            "--headless",
            "-env:UserInstallation=file:///tmp/lo_profile_vocaband",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pptx_path.parent),
            str(pptx_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    if not pdf_path.exists():
        raise RuntimeError(
            f"PDF conversion failed for {pptx_path}\n"
            f"stdout: {result.stdout}\nstderr: {result.stderr}"
        )
    return pdf_path


# ---------- Main ----------

def main():
    artefacts = []
    for lang in ("he", "ar"):
        s = STRINGS[lang]
        print(f"\n=== {s['tag']} ===")
        print("Building teacher handout…")
        artefacts.append(build_teacher_handout(s))
        print("Building manager handout…")
        artefacts.append(build_manager_handout(s))
        print("Building presentation deck…")
        artefacts.append(build_presentation(s))

    print("\nConverting to PDF…")
    for f in artefacts:
        print(f"  {f.name} -> PDF")
        to_pdf(f)

    print(f"\nDone. Outputs in {OUT_DIR}:")
    for f in sorted(OUT_DIR.iterdir()):
        size_kb = f.stat().st_size // 1024
        print(f"  {f.name}  ({size_kb} KB)")


if __name__ == "__main__":
    main()
